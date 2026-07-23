"""Text extraction for the file types StudyMate accepts.

Video is handled as a pre-existing transcript file (.srt/.vtt/.txt), not
audio-to-text — see README "What video transcripts means here".
"""

from pathlib import Path

import webvtt
from pptx import Presentation
from pypdf import PdfReader
from docx import Document as DocxDocument

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".srt", ".vtt", ".txt"}


class UnsupportedFileType(ValueError):
    pass


def doc_type_for(extension: str) -> str:
    if extension == ".pdf":
        return "pdf"
    if extension == ".docx":
        return "docx"
    if extension == ".pptx":
        return "pptx"
    if extension in (".srt", ".vtt"):
        return "transcript"
    return "text"


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".srt":
        return _extract_srt(path)
    if ext == ".vtt":
        return _extract_vtt(path)
    if ext == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")
    raise UnsupportedFileType(f"Unsupported file extension: {ext}")


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        texts.append(text)
        if texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


def _extract_srt(path: Path) -> str:
    # Minimal SRT parser: drop index lines and timecodes, keep spoken text.
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    text_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.isdigit():
            continue
        if "-->" in stripped:
            continue
        text_lines.append(stripped)
    return " ".join(text_lines)


def _extract_vtt(path: Path) -> str:
    captions = webvtt.read(str(path))
    return " ".join(c.text.replace("\n", " ") for c in captions)
